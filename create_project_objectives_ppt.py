from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for idx, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)


def add_arch_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "System Flow"

    left = Inches(0.7)
    top = Inches(1.4)
    width = Inches(2.5)
    height = Inches(1.1)

    boxes = [
        "Siren ML +\nLoRa/Wireless",
        "Dual Verification\nFusion",
        "Route + Hospital\nSelection",
        "Signal\nPreemption",
        "Hospital\nNotification",
    ]

    for i, text in enumerate(boxes):
        x = left + Inches(2.6) * i
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, width, height)
        shape.text = text
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(223, 240, 255)


def build_ppt(out_path: str) -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "Intelligent Emergency Traffic Preemption",
        "Hyderabad SUMO Project Objectives and Execution Plan",
    )

    add_bullet_slide(
        prs,
        "Primary Objectives",
        [
            "Detect emergency events using siren confidence and wireless/LoRa trigger",
            "Preempt upcoming traffic signals to build a green corridor",
            "Reroute ambulances dynamically to the best reachable hospital",
            "Notify destination hospitals with ETA and emergency type",
        ],
    )

    add_bullet_slide(
        prs,
        "Operational Targets",
        [
            "Run 10 to 15 ambulances in mixed high-density traffic",
            "Maintain visible emergency priority through signal preemption logs",
            "Guarantee hospital-edge reachability before simulation starts",
            "Restore traffic signal programs automatically after ambulance clearance",
        ],
    )

    add_bullet_slide(
        prs,
        "What Is Implemented",
        [
            "Map-style hospital symbols rendered in SUMO GUI",
            "Automatic traffic-light presence check and fallback net rebuild",
            "LoRa UDP listener option for wireless trigger ingestion",
            "Dynamic ETA-based rerouting and fair multi-ambulance conflict handling",
        ],
    )

    add_bullet_slide(
        prs,
        "How To Demonstrate",
        [
            "Run: .\\run_hyderabad.ps1 -Profile strict-production -SumoBinary sumo-gui",
            "Observe [SIGNAL] logs for phase switch and restore events",
            "Observe [INFO] route updates when hospital destination changes",
            "Optionally feed LoRa UDP payloads to trigger wireless confirmation",
        ],
    )

    add_arch_slide(prs)
    prs.save(out_path)


if __name__ == "__main__":
    build_ppt("Project_Objectives_Emergency_Traffic.pptx")
    print("Wrote Project_Objectives_Emergency_Traffic.pptx")
